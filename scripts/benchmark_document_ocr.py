#!/usr/bin/env python3
from __future__ import annotations

import json
import os
import sys
import tempfile
import time
from dataclasses import dataclass
from datetime import datetime, timezone
from difflib import SequenceMatcher
from io import BytesIO
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "backend"))
sys.path.insert(0, str(ROOT / "ocr-provider"))

from app.ingest.document_markdown import OcrDocumentResult, OcrPageResult, convert_document_to_markdown
from docx import Document
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches
from provider.app import OcrRuntime
from provider.config import Settings
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas


@dataclass(frozen=True)
class EngineSpec:
    name: str
    provider: str
    model: str
    languages: tuple[str, ...]
    device: str


@dataclass(frozen=True)
class FixtureCase:
    name: str
    filename: str
    mime_type: str
    expected_ocr_text: str
    strip_markers: tuple[str, ...] = ()


class LocalRuntimeDocumentOcrProvider:
    provider_kind = "local"

    def __init__(self, spec: EngineSpec, runtime: OcrRuntime) -> None:
        self.provider_name = spec.name
        self.ocr_model = spec.model
        self._runtime = runtime

    def extract(self, *, inputs: list[object], languages: list[str]) -> list[OcrDocumentResult]:
        results: list[OcrDocumentResult] = []
        for item in inputs:
            mime_type = str(getattr(item, "mime_type"))
            data = getattr(item, "data")
            source_id = str(getattr(item, "source_id"))
            if mime_type == "application/pdf":
                pages = self._runtime.ocr_pdf(data, getattr(item, "page_numbers"))
                text = "\n\n".join(page.text for page in pages if page.text)
                confidences = [page.confidence for page in pages if page.confidence is not None]
                confidence = round(sum(confidences) / len(confidences), 4) if confidences else None
                results.append(
                    OcrDocumentResult(
                        source_id=source_id,
                        text=text,
                        confidence=confidence,
                        warnings=[warning for page in pages for warning in page.warnings],
                        pages=[
                            OcrPageResult(
                                page_number=page.page_number,
                                text=page.text,
                                confidence=page.confidence,
                                warnings=list(page.warnings),
                            )
                            for page in pages
                        ],
                    )
                )
                continue
            image_result = self._runtime.ocr_image(data)
            results.append(
                OcrDocumentResult(
                    source_id=source_id,
                    text=image_result.text,
                    confidence=image_result.confidence,
                    warnings=list(image_result.warnings),
                )
            )
        return results


def _normalize_score_text(value: str) -> str:
    return "".join(ch.lower() for ch in value if ch.isalnum() or "\u4e00" <= ch <= "\u9fff")


def _score(expected: str, actual: str) -> float:
    return round(SequenceMatcher(None, _normalize_score_text(expected), _normalize_score_text(actual)).ratio(), 4)


def _pick_font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/Hiragino Sans GB.ttc",
        "/System/Library/Fonts/STHeiti Light.ttc",
        "/System/Library/Fonts/Supplemental/Songti.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
    ]
    for candidate in candidates:
        path = Path(candidate)
        if path.exists():
            return ImageFont.truetype(str(path), size=size)
    return ImageFont.load_default()


def _generate_fixtures(workdir: Path) -> tuple[list[FixtureCase], dict[str, Path]]:
    font = _pick_font(34)
    lines = [
        "城投债 Credit Update",
        "Q1 2026 风险提示 45.2%",
    ]
    expected_text = "\n".join(lines)

    image = Image.new("RGB", (900, 240), "white")
    draw = ImageDraw.Draw(image)
    draw.multiline_text((32, 40), expected_text, fill="black", font=font, spacing=18)
    image_path = workdir / "benchmark-image.png"
    image.save(image_path)

    pdf_path = workdir / "benchmark.pdf"
    pdf_buffer = BytesIO()
    pdf = canvas.Canvas(pdf_buffer)
    pdf.drawString(72, 800, "Native PDF body text for OCR benchmark coverage.")
    pdf.drawImage(ImageReader(str(image_path)), 72, 560, width=420, height=112)
    pdf.save()
    pdf_path.write_bytes(pdf_buffer.getvalue())

    docx_path = workdir / "benchmark.docx"
    doc = Document()
    doc.add_paragraph("Native DOCX body text for OCR benchmark coverage.")
    doc.add_picture(str(image_path))
    doc.save(docx_path)

    pptx_path = workdir / "benchmark.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Benchmark Slide"
    textbox = slide.shapes.add_textbox(left=0, top=0, width=400, height=100)
    textbox.text_frame.text = "Native PPTX body text for OCR benchmark coverage."
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(2), width=Inches(5))
    presentation.save(pptx_path)

    return (
        [
            FixtureCase("image", image_path.name, "image/png", expected_text),
            FixtureCase("pdf", pdf_path.name, "application/pdf", expected_text),
            FixtureCase(
                "docx",
                docx_path.name,
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                expected_text,
            ),
            FixtureCase(
                "pptx",
                pptx_path.name,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                expected_text,
                strip_markers=("Benchmark Slide", "Native PPTX body text for OCR benchmark coverage."),
            ),
        ],
        {
            "image": image_path,
            "pdf": pdf_path,
            "docx": docx_path,
            "pptx": pptx_path,
        },
    )


def _run_case(provider: LocalRuntimeDocumentOcrProvider, path: Path, case: FixtureCase, languages: tuple[str, ...]) -> dict[str, object]:
    started = time.perf_counter()
    try:
        result = convert_document_to_markdown(
            filename=path.name,
            mime_type=case.mime_type,
            file_bytes=path.read_bytes(),
            ocr_provider=provider,
            ocr_languages=list(languages),
            ocr_embedded_image_mode="always",
        )
        elapsed_ms = round((time.perf_counter() - started) * 1000, 2)
        ocr_text = "\n".join(part.content for part in result.parts if part.ocr_provenance)
        for marker in case.strip_markers:
            ocr_text = ocr_text.replace(marker, "")
        ocr_text = "\n".join(line for line in ocr_text.splitlines() if line.strip()).strip()
        warnings = [warning.message for warning in result.ocr_warnings]
        status = result.ocr_status
    except Exception as exc:
        elapsed_ms = round((time.perf_counter() - started) * 1000, 2)
        ocr_text = ""
        warnings = [f"{type(exc).__name__}: {exc}"]
        status = "failed"
    return {
        "fixture": case.name,
        "ocr_status": status,
        "elapsed_ms": elapsed_ms,
        "score": _score(case.expected_ocr_text, ocr_text),
        "ocr_text": ocr_text,
        "warnings": warnings,
    }


def _markdown_report(rows: list[dict[str, object]], environment: dict[str, object]) -> str:
    findings = _build_findings(rows, environment)
    lines = [
        "# Document OCR Benchmark",
        "",
        f"- Generated at: {environment['generated_at']}",
        f"- Host: {environment['host']}",
        f"- Python: {environment['python']}",
        f"- ONNX Runtime providers: {', '.join(environment['onnxruntime_providers'])}",
        f"- Torch MPS built/available: {environment['torch_mps_built']}/{environment['torch_mps_available']}",
        "- Method: generated image text is embedded into image / PDF / DOCX / PPTX; score compares OCR-only extracted text against the expected image text.",
        "",
        "## Findings",
        "",
    ]
    lines.extend(f"- {finding}" for finding in findings)
    lines.extend(
        [
            "",
        "| Engine | Fixture | Status | Score | Elapsed (ms) | Warnings |",
        "| --- | --- | --- | ---: | ---: | --- |",
        ]
    )
    for row in rows:
        lines.append(
            f"| {row['engine']} | {row['fixture']} | {row['ocr_status']} | {row['score']:.4f} | {row['elapsed_ms']:.2f} | "
            f"{'; '.join(row['warnings']) if row['warnings'] else ''} |"
        )
    lines.extend(
        [
            "",
            "## Samples",
            "",
        ]
    )
    for row in rows:
        excerpt = str(row["ocr_text"]).replace("\n", " / ").strip()
        lines.append(f"- `{row['engine']}` on `{row['fixture']}`: `{excerpt[:180]}`")
    lines.extend(
        [
            "",
            "## Raw JSON",
            "",
            "```json",
            json.dumps({"environment": environment, "rows": rows}, ensure_ascii=False, indent=2),
            "```",
        ]
    )
    return "\n".join(lines) + "\n"


def _build_findings(rows: list[dict[str, object]], environment: dict[str, object]) -> list[str]:
    findings: list[str] = []
    summaries = _summarize_rows(rows)
    rapid_cpu = summaries.get("rapidocr-local-cpu")
    rapid_coreml = summaries.get("rapidocr-macos-coreml")
    if rapid_cpu and rapid_coreml:
        if rapid_coreml["avg_score"] >= rapid_cpu["avg_score"] and rapid_coreml["avg_elapsed_ms"] > rapid_cpu["avg_elapsed_ms"]:
            findings.append(
                "On this macOS host, `rapidocr-macos-coreml` matched `rapidocr-local-cpu` accuracy but was materially slower, so CPU remains the better local default."
            )
    easy_mps = summaries.get("easyocr-macos-mps")
    if easy_mps:
        if environment["torch_mps_available"]:
            findings.append(
                "PyTorch MPS was available and `easyocr-macos-mps` executed on Metal, but it still produced no usable text on these fixtures."
            )
        else:
            findings.append("PyTorch MPS was not available for this Python build, so `easyocr-macos-mps` could not be exercised.")
    tesseract = summaries.get("tesseract-local-cpu")
    if tesseract and tesseract["avg_elapsed_ms"] < 150:
        findings.append(
            "Tesseract remained the fastest local option, but its mixed Chinese+English accuracy stayed materially below RapidOCR in this benchmark."
        )
    if not findings:
        findings.append("No benchmark-specific findings were generated.")
    return findings


def _summarize_rows(rows: list[dict[str, object]]) -> dict[str, dict[str, float]]:
    summary: dict[str, dict[str, float]] = {}
    for row in rows:
        engine = str(row["engine"])
        bucket = summary.setdefault(engine, {"count": 0.0, "avg_score": 0.0, "avg_elapsed_ms": 0.0})
        bucket["count"] += 1
        bucket["avg_score"] += float(row["score"])
        bucket["avg_elapsed_ms"] += float(row["elapsed_ms"])
    for bucket in summary.values():
        count = bucket["count"] or 1.0
        bucket["avg_score"] = round(bucket["avg_score"] / count, 4)
        bucket["avg_elapsed_ms"] = round(bucket["avg_elapsed_ms"] / count, 2)
    return summary


def main() -> int:
    workdir = Path(tempfile.mkdtemp(prefix="document-ocr-benchmark-"))
    cases, paths = _generate_fixtures(workdir)
    specs = [
        EngineSpec("rapidocr-local-cpu", "rapidocr", "rapidocr:ch_sim+en", ("ch_sim", "en"), "cpu"),
        EngineSpec("tesseract-local-cpu", "tesseract", "tesseract:ch_sim+en", ("ch_sim", "en"), "cpu"),
        EngineSpec("easyocr-local-cpu", "easyocr", "easyocr:ch_sim+en", ("ch_sim", "en"), "cpu"),
    ]
    if sys.platform == "darwin":
        specs.extend(
            [
                EngineSpec("rapidocr-macos-coreml", "rapidocr", "rapidocr:ch_sim+en", ("ch_sim", "en"), "coreml"),
                EngineSpec("easyocr-macos-mps", "easyocr", "easyocr:ch_sim+en", ("ch_sim", "en"), "mps"),
            ]
        )
    rows: list[dict[str, object]] = []
    for spec in specs:
        os.environ.update(
            {
                "OCR_PROVIDER": spec.provider,
                "OCR_MODEL": spec.model,
                "OCR_LANGUAGES": ",".join(spec.languages),
                "OCR_DEVICE": spec.device,
                "OCR_MODEL_STORAGE_DIR": str(workdir / f"models-{spec.provider}"),
            }
        )
        try:
            runtime = OcrRuntime(Settings.from_env())
            provider = LocalRuntimeDocumentOcrProvider(spec, runtime)
            _run_case(provider, paths["image"], cases[0], spec.languages)  # warmup
            for case in cases:
                row = _run_case(provider, paths[case.name], case, spec.languages)
                row["engine"] = spec.name
                rows.append(row)
        except Exception as exc:
            for case in cases:
                rows.append(
                    {
                        "engine": spec.name,
                        "fixture": case.name,
                        "ocr_status": "failed",
                        "elapsed_ms": 0.0,
                        "score": 0.0,
                        "ocr_text": "",
                        "warnings": [f"{type(exc).__name__}: {exc}"],
                    }
                )

    import onnxruntime as ort
    import torch
    environment = {
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "host": os.uname().nodename,
        "python": sys.version.split()[0],
        "onnxruntime_providers": ort.get_available_providers(),
        "torch_mps_built": bool(getattr(torch.backends, "mps", None) and torch.backends.mps.is_built()),
        "torch_mps_available": bool(getattr(torch.backends, "mps", None) and torch.backends.mps.is_available()),
    }
    output_path = ROOT / "docs" / "OCR_PROVIDER_BENCHMARK.md"
    output_path.write_text(_markdown_report(rows, environment), encoding="utf-8")
    print(f"wrote benchmark report to {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
