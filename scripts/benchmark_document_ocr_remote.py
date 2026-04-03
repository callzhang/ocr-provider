#!/usr/bin/env python3
from __future__ import annotations

import json
import os
import sys
import tempfile
import time
from dataclasses import dataclass
from datetime import datetime, timezone
from difflib import SequenceMatcher
from io import BytesIO
from pathlib import Path
from typing import Any

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "backend"))
sys.path.insert(0, str(ROOT / "ocr-provider"))

from app.ingest.document_markdown import (  # noqa: E402
    HttpDocumentOcrProvider,
    OcrDocumentResult,
    OcrInput as BackendOcrInput,
    OcrPageResult,
    convert_document_to_markdown,
)
from docx import Document  # noqa: E402
import fitz  # noqa: E402
from PIL import Image, ImageChops, ImageDraw, ImageEnhance, ImageFilter, ImageFont  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402
from provider.app import OcrRuntime  # noqa: E402
from provider.config import Settings  # noqa: E402
from reportlab.lib.utils import ImageReader  # noqa: E402
from reportlab.pdfgen import canvas  # noqa: E402


@dataclass(frozen=True)
class BenchmarkTarget:
    name: str
    kind: str
    base_url: str | None
    api_key: str | None
    model: str | None
    languages: tuple[str, ...]
    provider: str | None = None
    device: str = "cpu"
    model_storage_dir: Path | None = None
    render_scale: float = 2.0
    paragraph: bool = True
    tesseract_cmd: str | None = None
    paddle_text_detection_model_name: str | None = None
    paddle_text_recognition_model_name: str | None = None
    paddle_lang: str | None = None
    paddle_use_doc_orientation_classify: bool = False
    paddle_use_doc_unwarping: bool = False
    paddle_use_textline_orientation: bool = False
    paddle_disable_model_source_check: bool = True
    det_arch: str | None = None
    reco_arch: str | None = None


@dataclass(frozen=True)
class FixtureCase:
    name: str
    filename: str
    mime_type: str
    expected_ocr_text: str
    strip_markers: tuple[str, ...] = ()


@dataclass
class DirectDocumentOcrProvider:
    target: BenchmarkTarget

    def __post_init__(self) -> None:
        if not self.target.provider:
            raise ValueError(f"direct target {self.target.name} is missing provider")
        model_storage_dir = self.target.model_storage_dir or (ROOT / "runtime-cache" / self.target.name)
        self._settings = Settings(
            service_name=self.target.name,
            ocr_provider=self.target.provider,
            model_id=self.target.model or f"{self.target.provider}:{'+'.join(self.target.languages)}",
            model_alias=self.target.name,
            api_key=None,
            ocr_languages=self.target.languages,
            ocr_device=self.target.device,
            paragraph=self.target.paragraph,
            model_storage_dir=model_storage_dir,
            render_scale=self.target.render_scale,
            tesseract_cmd=self.target.tesseract_cmd,
            paddle_text_detection_model_name=self.target.paddle_text_detection_model_name,
            paddle_text_recognition_model_name=self.target.paddle_text_recognition_model_name,
            paddle_lang=self.target.paddle_lang,
            paddle_use_doc_orientation_classify=self.target.paddle_use_doc_orientation_classify,
            paddle_use_doc_unwarping=self.target.paddle_use_doc_unwarping,
            paddle_use_textline_orientation=self.target.paddle_use_textline_orientation,
            paddle_disable_model_source_check=self.target.paddle_disable_model_source_check,
        )
        self._runtime = OcrRuntime(self._settings)
        self.provider_kind = "direct"
        self.provider_name = self.target.name
        self.ocr_model = self._settings.model_id

    def extract(self, *, inputs: list[BackendOcrInput], languages: list[str]) -> list[OcrDocumentResult]:
        _ = languages
        results: list[OcrDocumentResult] = []
        for item in inputs:
            if item.mime_type.lower() == "application/pdf":
                pages = self._runtime.ocr_pdf(item.data, item.page_numbers)
                joined = "\n\n".join(page.text for page in pages if page.text)
                confidences = [page.confidence for page in pages if page.confidence is not None]
                confidence = round(sum(confidences) / len(confidences), 4) if confidences else None
                results.append(
                    OcrDocumentResult(
                        source_id=item.source_id,
                        text=joined,
                        confidence=confidence,
                        warnings=[warning for page in pages for warning in page.warnings],
                        pages=[
                            OcrPageResult(
                                page_number=page.page_number,
                                text=page.text,
                                confidence=page.confidence,
                                warnings=list(page.warnings),
                            )
                            for page in pages
                        ],
                    )
                )
                continue
            image_result = self._runtime.ocr_image(item.data)
            results.append(
                OcrDocumentResult(
                    source_id=item.source_id,
                    text=image_result.text,
                    confidence=image_result.confidence,
                    warnings=list(image_result.warnings),
                    pages=[],
                )
            )
        return results


@dataclass
class SuryaDocumentOcrProvider:
    target: BenchmarkTarget

    def __post_init__(self) -> None:
        cache_root = self.target.model_storage_dir or (ROOT / "runtime-cache" / self.target.name)
        cache_root.mkdir(parents=True, exist_ok=True)
        os.environ.setdefault("XDG_CACHE_HOME", str(cache_root))
        os.environ.setdefault("HF_HOME", str(cache_root / "hf-home"))
        os.environ.setdefault("HF_HUB_CACHE", str(cache_root / "hf-hub"))
        os.environ.setdefault("TORCH_HOME", str(cache_root / "torch-home"))

        from surya.detection import DetectionPredictor
        from surya.foundation import FoundationPredictor
        from surya.recognition import RecognitionPredictor

        device = "mps" if self.target.device in {"mps", "coreml"} else self.target.device
        self._render_scale = self.target.render_scale
        self._foundation = FoundationPredictor(device=device)
        self._detector = DetectionPredictor(device=device)
        self._recognizer = RecognitionPredictor(self._foundation)
        self.provider_kind = "direct"
        self.provider_name = self.target.name
        self.ocr_model = self.target.model or "surya-ocr"

    def extract(self, *, inputs: list[BackendOcrInput], languages: list[str]) -> list[OcrDocumentResult]:
        _ = languages
        results: list[OcrDocumentResult] = []
        for item in inputs:
            if item.mime_type.lower() == "application/pdf":
                results.append(self._ocr_pdf(item))
                continue
            image_result = self._ocr_image(item.data)
            results.append(
                OcrDocumentResult(
                    source_id=item.source_id,
                    text=image_result.text,
                    confidence=image_result.confidence,
                    warnings=list(image_result.warnings),
                )
            )
        return results

    def _ocr_pdf(self, item: BackendOcrInput) -> OcrDocumentResult:
        pdf = fitz.open(stream=item.data, filetype="pdf")
        requested = item.page_numbers or list(range(1, pdf.page_count + 1))
        matrix = fitz.Matrix(self._render_scale, self._render_scale)
        page_results: list[OcrPageResult] = []
        for page_number in requested:
            if page_number < 1 or page_number > pdf.page_count:
                page_results.append(
                    OcrPageResult(
                        page_number=page_number,
                        text="",
                        warnings=[f"page {page_number} is out of range for PDF with {pdf.page_count} pages"],
                    )
                )
                continue
            page = pdf.load_page(page_number - 1)
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            result = self._ocr_image(pixmap.tobytes("png"))
            page_results.append(
                OcrPageResult(
                    page_number=page_number,
                    text=result.text,
                    confidence=result.confidence,
                    warnings=list(result.warnings) or ([] if result.text else [f"no text recognized on page {page_number}"]),
                )
            )
        confidences = [page.confidence for page in page_results if page.confidence is not None]
        return OcrDocumentResult(
            source_id=item.source_id,
            text="\n\n".join(page.text for page in page_results if page.text),
            confidence=round(sum(confidences) / len(confidences), 4) if confidences else None,
            warnings=[warning for page in page_results for warning in page.warnings],
            pages=page_results,
        )

    def _ocr_image(self, data: bytes):
        image = Image.open(BytesIO(data)).convert("RGB")
        predictions = self._recognizer(
            [image],
            det_predictor=self._detector,
            highres_images=[image],
            math_mode=False,
            sort_lines=True,
            return_words=True,
        )
        prediction = predictions[0] if predictions else None
        if prediction is None:
            return type("SuryaResult", (), {"text": "", "confidence": None, "warnings": ["surya returned no predictions"]})()
        lines = [line.text.strip() for line in prediction.text_lines if line.text and line.text.strip()]
        confidences = [
            float(line.confidence)
            for line in prediction.text_lines
            if getattr(line, "confidence", None) is not None
        ]
        confidence = round(sum(confidences) / len(confidences), 4) if confidences else None
        warnings = [] if lines else ["surya returned no text lines"]
        return type("SuryaResult", (), {"text": "\n".join(lines).strip(), "confidence": confidence, "warnings": warnings})()


@dataclass
class OnnxtrDocumentOcrProvider:
    target: BenchmarkTarget

    def __post_init__(self) -> None:
        cache_root = self.target.model_storage_dir or (ROOT / "runtime-cache" / self.target.name)
        cache_root.mkdir(parents=True, exist_ok=True)
        os.environ.setdefault("ONNXTR_CACHE_DIR", str(cache_root))

        from onnxtr.models import EngineConfig, ocr_predictor

        self._providers = self._resolve_providers()
        self._engine_cfg = EngineConfig(providers=self._providers)
        self._predictor = ocr_predictor(
            det_arch=self.target.det_arch or "fast_base",
            reco_arch=self.target.reco_arch or "crnn_vgg16_bn",
            det_engine_cfg=self._engine_cfg,
            reco_engine_cfg=self._engine_cfg,
            clf_engine_cfg=self._engine_cfg,
        )
        self.provider_kind = "direct"
        self.provider_name = self.target.name
        self.ocr_model = self.target.model or f"onnxtr:{self.target.det_arch or 'fast_base'}+{self.target.reco_arch or 'crnn_vgg16_bn'}"

    def _resolve_providers(self) -> list[tuple[str, dict[str, object]]]:
        normalized = self.target.device.lower()
        cpu = ("CPUExecutionProvider", {"arena_extend_strategy": "kSameAsRequested"})
        if normalized in {"mps", "coreml"}:
            return [("CoreMLExecutionProvider", {}), cpu]
        return [cpu]

    def extract(self, *, inputs: list[BackendOcrInput], languages: list[str]) -> list[OcrDocumentResult]:
        _ = languages
        results: list[OcrDocumentResult] = []
        for item in inputs:
            if item.mime_type.lower() == "application/pdf":
                results.append(self._ocr_pdf(item))
                continue
            text, confidence, warnings = self._ocr_bytes(item.data, item.mime_type)
            results.append(
                OcrDocumentResult(
                    source_id=item.source_id,
                    text=text,
                    confidence=confidence,
                    warnings=warnings,
                )
            )
        return results

    def _ocr_pdf(self, item: BackendOcrInput) -> OcrDocumentResult:
        from onnxtr.io import DocumentFile

        requested_pages = item.page_numbers or []
        document = DocumentFile.from_pdf(BytesIO(item.data))
        raw_result = self._predictor(document)
        pages = list(raw_result.pages)
        if requested_pages:
            page_indexes = {page_number - 1 for page_number in requested_pages if page_number >= 1}
            filtered_pages = [page for index, page in enumerate(pages) if index in page_indexes]
        else:
            filtered_pages = pages
        page_results: list[OcrPageResult] = []
        for index, page in enumerate(filtered_pages, start=1):
            page_text = page.render().strip()
            page_confidence = _average(_collect_onnxtr_confidences(page))
            page_results.append(
                OcrPageResult(
                    page_number=requested_pages[index - 1] if requested_pages else index,
                    text=page_text,
                    confidence=page_confidence,
                    warnings=[] if page_text else [f"onnxtr returned no text for page {index}"],
                )
            )
        confidences = [page.confidence for page in page_results if page.confidence is not None]
        return OcrDocumentResult(
            source_id=item.source_id,
            text="\n\n".join(page.text for page in page_results if page.text),
            confidence=round(sum(confidences) / len(confidences), 4) if confidences else None,
            warnings=[warning for page in page_results for warning in page.warnings],
            pages=page_results,
        )

    def _ocr_bytes(self, data: bytes, mime_type: str) -> tuple[str, float | None, list[str]]:
        from onnxtr.io import DocumentFile

        if mime_type.startswith("image/"):
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as handle:
                handle.write(data)
                temp_path = Path(handle.name)
            try:
                document = DocumentFile.from_images([str(temp_path)])
                raw_result = self._predictor(document)
            finally:
                temp_path.unlink(missing_ok=True)
        else:
            raise ValueError(f"unsupported mime_type for onnxtr benchmark provider: {mime_type}")
        text = raw_result.render().strip()
        confidence = _average(_collect_onnxtr_confidences(raw_result))
        warnings = [] if text else ["onnxtr returned no text"]
        return text, confidence, warnings


def _pick_font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/Hiragino Sans GB.ttc",
        "/System/Library/Fonts/Supplemental/Songti.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size=size)
    return ImageFont.load_default()


def _draw_multiline_image(
    text: str,
    *,
    size: tuple[int, int],
    font_size: int,
    fill: str,
    background: str,
    spacing: int,
) -> Image.Image:
    image = Image.new("RGB", size, background)
    draw = ImageDraw.Draw(image)
    draw.multiline_text((24, 28), text, fill=fill, font=_pick_font(font_size), spacing=spacing)
    return image


def _draw_dense_column_image() -> tuple[Image.Image, str]:
    lines = [
        "2026年城投平台再融资跟踪",
        "Left Column",
        "1. 净融资 -2.3bn",
        "2. 非标压降 17.4%",
        "3. 平台现金短债比 1.18x",
        "Right Column",
        "A. 债项评级维持 AA+",
        "B. 土储去化周期 19.6月",
        "C. Watchlist: 贵州、云南",
        "Footnote: EBITDA / Interest 2.7x",
    ]
    image = Image.new("RGB", (1320, 760), "#fbfcfe")
    draw = ImageDraw.Draw(image)
    font_title = _pick_font(34)
    font_body = _pick_font(24)
    font_note = _pick_font(20)
    draw.rounded_rectangle((26, 24, 1294, 734), radius=16, outline="#c5cfdb", width=3)
    draw.text((52, 46), lines[0], fill="#111827", font=font_title)
    draw.line((650, 112, 650, 700), fill="#d6dde6", width=3)
    left_block = "\n".join(lines[1:5])
    right_block = "\n".join(lines[5:9])
    draw.multiline_text((58, 132), left_block, fill="#172033", font=font_body, spacing=22)
    draw.multiline_text((690, 132), right_block, fill="#172033", font=font_body, spacing=22)
    draw.text((58, 680), lines[9], fill="#5b6472", font=font_note)
    expected = "\n".join(lines)
    return image, expected


def _generate_fixtures(workdir: Path) -> tuple[list[FixtureCase], dict[str, Path]]:
    clean_text = "城投债 Credit Update\nQ1 2026 风险提示 45.2%"
    clean_image = _draw_multiline_image(clean_text, size=(900, 240), font_size=34, fill="black", background="white", spacing=18)
    clean_path = workdir / "challenge-clean.png"
    clean_image.save(clean_path)

    low_contrast_text = "并表口径 Revenue 1.28bn\n现金回收率 93.4% / watchlist"
    low_contrast = _draw_multiline_image(low_contrast_text, size=(960, 260), font_size=24, fill="#59606b", background="#eef1f5", spacing=14)
    low_contrast = ImageEnhance.Contrast(low_contrast).enhance(0.92).filter(ImageFilter.GaussianBlur(radius=0.4))
    low_contrast_path = workdir / "challenge-low-contrast.png"
    low_contrast.save(low_contrast_path)

    rotated_text = "2026Q1 到期分布 8.2bn\nAA+ / 城投平台 / 债务久期 3.4y"
    rotated_base = _draw_multiline_image(rotated_text, size=(960, 260), font_size=26, fill="black", background="white", spacing=16)
    rotated = rotated_base.rotate(7.2, expand=True, fillcolor="white")
    rotated = ImageChops.offset(rotated, 10, 6)
    rotated_path = workdir / "challenge-rotated.png"
    rotated.save(rotated_path)

    table_text = "项目 数值\n票息 4.52%\n主体评级 AA+\n债券余额 38.6bn"
    table = Image.new("RGB", (980, 320), "white")
    draw = ImageDraw.Draw(table)
    font = _pick_font(24)
    draw.rounded_rectangle((18, 18, 962, 302), radius=12, outline="#8b95a7", width=2)
    for y in (78, 138, 198, 258):
        draw.line((18, y, 962, y), fill="#c7ced9", width=2)
    draw.line((220, 18, 220, 302), fill="#c7ced9", width=2)
    rows = [("项目", "数值"), ("票息", "4.52%"), ("主体评级", "AA+"), ("债券余额", "38.6bn")]
    for idx, (left, right) in enumerate(rows):
        y = 34 + idx * 60
        draw.text((42, y), left, fill="#1f2937", font=font)
        draw.text((256, y), right, fill="#111827", font=font)
    table_path = workdir / "challenge-table.png"
    table.save(table_path)

    dense_column_image, dense_column_text = _draw_dense_column_image()
    dense_column = ImageEnhance.Contrast(dense_column_image).enhance(0.97).filter(ImageFilter.GaussianBlur(radius=0.25))
    dense_column_path = workdir / "challenge-dense-columns.png"
    dense_column.save(dense_column_path)

    combined_text = "\n".join([clean_text, low_contrast_text, rotated_text, table_text, dense_column_text])
    pdf_path = workdir / "challenge-doc.pdf"
    pdf_buffer = BytesIO()
    pdf = canvas.Canvas(pdf_buffer)
    pdf.drawString(72, 810, "Remote OCR benchmark document")
    pdf.drawImage(ImageReader(str(clean_path)), 72, 640, width=420, height=112)
    pdf.drawImage(ImageReader(str(low_contrast_path)), 72, 430, width=450, height=122)
    pdf.drawImage(ImageReader(str(rotated_path)), 72, 220, width=450, height=122)
    pdf.showPage()
    pdf.drawImage(ImageReader(str(table_path)), 72, 470, width=470, height=153)
    pdf.drawImage(ImageReader(str(dense_column_path)), 48, 84, width=520, height=299)
    pdf.save()
    pdf_path.write_bytes(pdf_buffer.getvalue())

    docx_path = workdir / "challenge-doc.docx"
    doc = Document()
    doc.add_paragraph("Remote OCR benchmark document.")
    for image_path in (clean_path, low_contrast_path, rotated_path, table_path, dense_column_path):
        doc.add_picture(str(image_path))
    doc.save(docx_path)

    pptx_path = workdir / "challenge-doc.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Remote OCR Challenge"
    slide.shapes.add_picture(str(clean_path), Inches(0.6), Inches(1.0), width=Inches(4.8))
    slide.shapes.add_picture(str(low_contrast_path), Inches(0.6), Inches(3.0), width=Inches(5.0))
    slide2 = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide2.shapes.title.text = "Remote OCR Challenge 2"
    slide2.shapes.add_picture(str(rotated_path), Inches(0.6), Inches(1.0), width=Inches(5.0))
    slide2.shapes.add_picture(str(table_path), Inches(0.6), Inches(3.1), width=Inches(5.1))
    slide3 = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide3.shapes.title.text = "Remote OCR Challenge 3"
    slide3.shapes.add_picture(str(dense_column_path), Inches(0.4), Inches(1.2), width=Inches(8.5))
    presentation.save(pptx_path)

    return (
        [
            FixtureCase("image-clean", clean_path.name, "image/png", clean_text),
            FixtureCase("image-low-contrast", low_contrast_path.name, "image/png", low_contrast_text),
            FixtureCase("image-rotated", rotated_path.name, "image/png", rotated_text),
            FixtureCase("image-table", table_path.name, "image/png", table_text),
            FixtureCase("image-dense-columns", dense_column_path.name, "image/png", dense_column_text),
            FixtureCase("pdf-challenge", pdf_path.name, "application/pdf", combined_text),
            FixtureCase(
                "docx-challenge",
                docx_path.name,
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                combined_text,
            ),
            FixtureCase(
                "pptx-challenge",
                pptx_path.name,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                combined_text,
                strip_markers=("Remote OCR Challenge", "Remote OCR Challenge 2", "Remote OCR Challenge 3"),
            ),
        ],
        {
            "image-clean": clean_path,
            "image-low-contrast": low_contrast_path,
            "image-rotated": rotated_path,
            "image-table": table_path,
            "image-dense-columns": dense_column_path,
            "pdf-challenge": pdf_path,
            "docx-challenge": docx_path,
            "pptx-challenge": pptx_path,
        },
    )


def _normalize_score_text(value: str) -> str:
    return "".join(ch.lower() for ch in value if ch.isalnum() or "\u4e00" <= ch <= "\u9fff")


def _score(expected: str, actual: str) -> float:
    return round(SequenceMatcher(None, _normalize_score_text(expected), _normalize_score_text(actual)).ratio(), 4)


def _build_provider(target: BenchmarkTarget):
    if target.kind == "surya":
        return SuryaDocumentOcrProvider(target)
    if target.kind == "onnxtr":
        return OnnxtrDocumentOcrProvider(target)
    if target.kind == "direct":
        return DirectDocumentOcrProvider(target)
    if not target.base_url:
        raise ValueError(f"http target {target.name} is missing base_url")
    return HttpDocumentOcrProvider(
        base_url=target.base_url,
        api_key=target.api_key,
        model=target.model,
        timeout_seconds=180.0,
        provider_name=target.name,
    )


def _run_case(target: BenchmarkTarget, provider: Any, path: Path, case: FixtureCase) -> dict[str, object]:
    started = time.perf_counter()
    try:
        result = convert_document_to_markdown(
            filename=path.name,
            mime_type=case.mime_type,
            file_bytes=path.read_bytes(),
            ocr_provider=provider,
            ocr_languages=list(target.languages),
            ocr_embedded_image_mode="always",
        )
        elapsed_ms = round((time.perf_counter() - started) * 1000, 2)
        ocr_text = "\n".join(part.content for part in result.parts if part.ocr_provenance)
        for marker in case.strip_markers:
            ocr_text = ocr_text.replace(marker, "")
        ocr_text = "\n".join(line for line in ocr_text.splitlines() if line.strip()).strip()
        warnings = [warning.message for warning in result.ocr_warnings]
        status = result.ocr_status
    except Exception as exc:
        elapsed_ms = round((time.perf_counter() - started) * 1000, 2)
        ocr_text = ""
        warnings = [f"{type(exc).__name__}: {exc}"]
        status = "failed"
    return {
        "engine": target.name,
        "kind": target.kind,
        "fixture": case.name,
        "ocr_status": status,
        "elapsed_ms": elapsed_ms,
        "score": _score(case.expected_ocr_text, ocr_text),
        "ocr_text": ocr_text,
        "warnings": warnings,
    }


def _markdown_report(rows: list[dict[str, object]], environment: dict[str, object]) -> str:
    by_engine: dict[str, list[dict[str, object]]] = {}
    for row in rows:
        by_engine.setdefault(str(row["engine"]), []).append(row)

    lines = [
        "# Accelerated OCR Provider Benchmark",
        "",
        f"- Generated at: {environment['generated_at']}",
        f"- Host: {environment['host']}",
        f"- Python: {environment['python']}",
        "- Method: harder synthetic challenge set with low contrast, rotation, table screenshots, and a dense two-column screenshot; image cases are tested directly, PDF/DOCX/PPTX cases embed multiple challenge images and run through the normal document markdown ingest path.",
        "",
        "## Engine Summary",
        "",
        "| Engine | Kind | Success | Avg Score | Avg Elapsed (ms) |",
        "| --- | --- | ---: | ---: | ---: |",
    ]
    for engine, engine_rows in sorted(by_engine.items()):
        avg_score = round(sum(float(item["score"]) for item in engine_rows) / len(engine_rows), 4)
        avg_elapsed = round(sum(float(item["elapsed_ms"]) for item in engine_rows) / len(engine_rows), 2)
        success_count = sum(1 for item in engine_rows if str(item["ocr_status"]) == "succeeded")
        lines.append(f"| {engine} | {engine_rows[0]['kind']} | {success_count}/{len(engine_rows)} | {avg_score:.4f} | {avg_elapsed:.2f} |")

    lines.extend(
        [
            "",
            "## Per Fixture",
            "",
            "| Engine | Fixture | Status | Score | Elapsed (ms) | Warnings |",
            "| --- | --- | --- | ---: | ---: | --- |",
        ]
    )
    for row in rows:
        lines.append(
            f"| {row['engine']} | {row['fixture']} | {row['ocr_status']} | {row['score']:.4f} | {row['elapsed_ms']:.2f} | "
            f"{'; '.join(row['warnings']) if row['warnings'] else ''} |"
        )

    lines.extend(["", "## Samples", ""])
    for row in rows:
        lines.append(f"- `{row['engine']}` on `{row['fixture']}`: `{str(row['ocr_text']).replace(chr(10), ' / ')[:220]}`")
    lines.extend(["", "## Raw JSON", "", "```json", json.dumps({"environment": environment, "rows": rows}, ensure_ascii=False, indent=2), "```"])
    return "\n".join(lines) + "\n"


def _load_targets() -> list[BenchmarkTarget]:
    raw = os.environ.get("REMOTE_OCR_TARGETS_JSON", "").strip()
    if not raw:
        raise SystemExit("REMOTE_OCR_TARGETS_JSON is required")
    parsed = json.loads(raw)
    targets: list[BenchmarkTarget] = []
    for item in parsed:
        api_key = item.get("api_key")
        api_key_env = item.get("api_key_env")
        if api_key_env:
            api_key = os.environ.get(str(api_key_env))
        targets.append(
            BenchmarkTarget(
                name=str(item["name"]),
                kind=str(item.get("kind") or "http"),
                base_url=str(item["base_url"]) if item.get("base_url") else None,
                api_key=api_key,
                model=item.get("model"),
                languages=tuple(item.get("languages") or ["ch_sim", "en"]),
                provider=item.get("provider"),
                device=str(item.get("device") or "cpu"),
                model_storage_dir=Path(item["model_storage_dir"]) if item.get("model_storage_dir") else None,
                render_scale=float(item.get("render_scale") or 2.0),
                paragraph=bool(item.get("paragraph", True)),
                tesseract_cmd=item.get("tesseract_cmd"),
                paddle_text_detection_model_name=item.get("paddle_text_detection_model_name"),
                paddle_text_recognition_model_name=item.get("paddle_text_recognition_model_name"),
                paddle_lang=item.get("paddle_lang"),
                paddle_use_doc_orientation_classify=bool(item.get("paddle_use_doc_orientation_classify", False)),
                paddle_use_doc_unwarping=bool(item.get("paddle_use_doc_unwarping", False)),
                paddle_use_textline_orientation=bool(item.get("paddle_use_textline_orientation", False)),
                paddle_disable_model_source_check=bool(item.get("paddle_disable_model_source_check", True)),
                det_arch=item.get("det_arch"),
                reco_arch=item.get("reco_arch"),
            )
        )
    return targets


def _average(values: list[float]) -> float | None:
    return round(sum(values) / len(values), 4) if values else None


def _collect_onnxtr_confidences(element: Any) -> list[float]:
    confidences: list[float] = []
    if hasattr(element, "pages"):
        for page in getattr(element, "pages", []):
            confidences.extend(_collect_onnxtr_confidences(page))
        return confidences
    if hasattr(element, "blocks"):
        for block in getattr(element, "blocks", []):
            confidences.extend(_collect_onnxtr_confidences(block))
        return confidences
    if hasattr(element, "lines"):
        for line in getattr(element, "lines", []):
            confidences.extend(_collect_onnxtr_confidences(line))
        return confidences
    if hasattr(element, "words"):
        for word in getattr(element, "words", []):
            value = getattr(word, "confidence", None)
            if value is None:
                continue
            try:
                confidences.append(float(value))
            except (TypeError, ValueError):
                continue
    return confidences


def main() -> int:
    targets = _load_targets()
    workdir = Path(tempfile.mkdtemp(prefix="document-ocr-remote-benchmark-"))
    cases, paths = _generate_fixtures(workdir)
    rows: list[dict[str, object]] = []
    for target in targets:
        provider = _build_provider(target)
        _run_case(target, provider, paths[cases[0].name], cases[0])
        for case in cases:
            rows.append(_run_case(target, provider, paths[case.name], case))

    environment = {
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "host": os.uname().nodename,
        "python": sys.version.split()[0],
    }
    output_env = os.environ.get("OCR_BENCHMARK_OUTPUT_PATH", "").strip()
    output_path = Path(output_env) if output_env else (ROOT / "docs" / "OCR_PROVIDER_REMOTE_BENCHMARK.md")
    output_path.write_text(_markdown_report(rows, environment), encoding="utf-8")
    print(f"wrote benchmark report to {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
