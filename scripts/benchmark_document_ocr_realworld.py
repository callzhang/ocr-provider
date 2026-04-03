#!/usr/bin/env python3
"""Real-world OCR benchmark — image-based fixtures, pypdf ground truth.

All four format containers (image, PDF, DOCX, PPTX) embed the same rendered
page from a real Chinese financial analyst report (华源证券·中国平安研报).
Ground truth is extracted from the original native-text PDF via pypdf, giving
a clean baseline that reflects what perfect OCR should recover.

Score = SequenceMatcher ratio against pypdf ground truth (0–1).
"""
from __future__ import annotations

import json
import os
import re
import sys
import tempfile
import time
from dataclasses import dataclass
from datetime import datetime, timezone
from difflib import SequenceMatcher
from io import BytesIO
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
FIXTURES_DIR = ROOT / "scripts" / "realworld-fixtures"
sys.path.insert(0, str(ROOT / "backend"))
sys.path.insert(0, str(ROOT / "ocr-provider"))

from app.ingest.document_markdown import OcrDocumentResult, OcrPageResult, convert_document_to_markdown
from docx import Document
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from provider.app import OcrRuntime
from provider.config import Settings
from pypdf import PdfReader
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas as rl_canvas


# ---------------------------------------------------------------------------
# Source files
# ---------------------------------------------------------------------------
ANALYST_PDF = FIXTURES_DIR / "dfcfw-pingan-analyst.pdf"
ANALYST_PNG = FIXTURES_DIR / "dfcfw-page1.png"   # pre-rendered at 150 dpi


# ---------------------------------------------------------------------------
# Ground truth
# ---------------------------------------------------------------------------

def _extract_ground_truth(pdf_path: Path, page_index: int = 0) -> str:
    """Extract native text from a PDF page via pypdf and clean it up."""
    reader = PdfReader(str(pdf_path))
    raw = reader.pages[page_index].extract_text() or ""
    # Remove known PDF template artifacts
    raw = raw.replace("hyzqdatemark", "")
    # Collapse excessive blank lines
    raw = re.sub(r"\n{3,}", "\n\n", raw).strip()
    return raw


# ---------------------------------------------------------------------------
# Scoring
# ---------------------------------------------------------------------------

def _normalize(text: str) -> str:
    """Keep only alphanumeric + CJK for comparison."""
    return "".join(
        ch.lower() for ch in text
        if ch.isalnum() or "\u4e00" <= ch <= "\u9fff"
    )


def _score(expected: str, actual: str) -> float:
    return round(
        SequenceMatcher(None, _normalize(expected), _normalize(actual)).ratio(), 4
    )


# ---------------------------------------------------------------------------
# Fixture builders
# ---------------------------------------------------------------------------

def _image_in_pdf(png_bytes: bytes, dpi: int = 150) -> bytes:
    """Wrap a PNG as a single-page image-only PDF."""
    img = Image.open(BytesIO(png_bytes))
    w_pt = img.width * 72 / dpi
    h_pt = img.height * 72 / dpi
    buf = BytesIO()
    c = rl_canvas.Canvas(buf, pagesize=(w_pt, h_pt))
    c.drawImage(ImageReader(BytesIO(png_bytes)), 0, 0, width=w_pt, height=h_pt)
    c.save()
    return buf.getvalue()


def _image_in_docx(png_bytes: bytes) -> bytes:
    doc = Document()
    doc.add_picture(BytesIO(png_bytes))
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _image_in_pptx(png_bytes: bytes) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    slide.shapes.add_picture(BytesIO(png_bytes), Inches(0), Inches(0), width=Inches(10))
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


@dataclass(frozen=True)
class Fixture:
    name: str
    filename: str
    mime_type: str
    file_bytes: bytes


def _build_fixtures(png_bytes: bytes) -> list[Fixture]:
    return [
        Fixture(
            name="image",
            filename="analyst-report-page1.png",
            mime_type="image/png",
            file_bytes=png_bytes,
        ),
        Fixture(
            name="pdf",
            filename="analyst-report-page1.pdf",
            mime_type="application/pdf",
            file_bytes=_image_in_pdf(png_bytes),
        ),
        Fixture(
            name="docx",
            filename="analyst-report-page1.docx",
            mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            file_bytes=_image_in_docx(png_bytes),
        ),
        Fixture(
            name="pptx",
            filename="analyst-report-page1.pptx",
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            file_bytes=_image_in_pptx(png_bytes),
        ),
    ]


# ---------------------------------------------------------------------------
# Engine wrapper
# ---------------------------------------------------------------------------

@dataclass(frozen=True)
class EngineSpec:
    name: str
    provider: str
    model: str
    languages: tuple[str, ...]
    device: str


class LocalRuntimeOcrProvider:
    provider_kind = "local"

    def __init__(self, spec: EngineSpec, runtime: OcrRuntime) -> None:
        self.provider_name = spec.name
        self.ocr_model = spec.model
        self._runtime = runtime

    def extract(self, *, inputs: list[object], languages: list[str]) -> list[OcrDocumentResult]:
        results: list[OcrDocumentResult] = []
        for item in inputs:
            mime_type = str(getattr(item, "mime_type"))
            data = getattr(item, "data")
            source_id = str(getattr(item, "source_id"))
            if mime_type == "application/pdf":
                pages = self._runtime.ocr_pdf(data, getattr(item, "page_numbers"))
                text = "\n\n".join(p.text for p in pages if p.text)
                confidences = [p.confidence for p in pages if p.confidence is not None]
                confidence = round(sum(confidences) / len(confidences), 4) if confidences else None
                results.append(OcrDocumentResult(
                    source_id=source_id, text=text, confidence=confidence,
                    warnings=[w for p in pages for w in p.warnings],
                    pages=[OcrPageResult(page_number=p.page_number, text=p.text,
                                        confidence=p.confidence, warnings=list(p.warnings))
                           for p in pages],
                ))
            else:
                r = self._runtime.ocr_image(data)
                results.append(OcrDocumentResult(
                    source_id=source_id, text=r.text,
                    confidence=r.confidence, warnings=list(r.warnings),
                ))
        return results


# ---------------------------------------------------------------------------
# Run one fixture
# ---------------------------------------------------------------------------

def _run(provider: LocalRuntimeOcrProvider, fixture: Fixture,
         languages: tuple[str, ...], ground_truth: str) -> dict:
    started = time.perf_counter()
    try:
        result = convert_document_to_markdown(
            filename=fixture.filename,
            mime_type=fixture.mime_type,
            file_bytes=fixture.file_bytes,
            ocr_provider=provider,
            ocr_languages=list(languages),
            ocr_embedded_image_mode="always",
        )
        elapsed_ms = round((time.perf_counter() - started) * 1000, 2)
        ocr_text = "\n".join(
            part.content for part in result.parts if part.ocr_provenance
        ).strip()
        warnings = [w.message for w in result.ocr_warnings]
        status = result.ocr_status
    except Exception as exc:
        elapsed_ms = round((time.perf_counter() - started) * 1000, 2)
        ocr_text = ""
        warnings = [f"{type(exc).__name__}: {exc}"]
        status = "failed"

    return {
        "fixture": fixture.name,
        "ocr_status": status,
        "elapsed_ms": elapsed_ms,
        "score": _score(ground_truth, ocr_text),
        "ocr_text": ocr_text,
        "warnings": warnings,
    }


# ---------------------------------------------------------------------------
# Report
# ---------------------------------------------------------------------------

def _markdown_report(ground_truth: str, rows: list[dict], env: dict) -> str:
    lines = [
        "# Real-World OCR Benchmark",
        "",
        f"- Generated at: {env['generated_at']}",
        f"- Host: {env['host']}",
        f"- Python: {env['python']}",
        f"- ONNX Runtime providers: {', '.join(env['onnxruntime_providers'])}",
        "",
        "## Source & Ground Truth",
        "",
        "- **Image source**: 华源证券·中国平安研究报告 page 1",
        "  rendered from `dfcfw-pingan-analyst.pdf` at 150 dpi",
        "- **Ground truth**: pypdf native-text extraction from the same PDF page",
        "- **Score**: SequenceMatcher ratio (0–1) of normalized OCR text vs. ground truth",
        "- **All 4 fixtures** use the same rendered PNG — only the container format differs",
        "",
        "### Ground Truth Excerpt (first 400 chars)",
        "",
        "```",
        ground_truth[:400],
        "```",
        "",
        "## Results",
        "",
        "| Engine | Fixture | Status | Score | Elapsed (ms) | Warnings |",
        "| --- | --- | --- | ---: | ---: | --- |",
    ]
    for row in rows:
        warn = "; ".join(row["warnings"]) if row["warnings"] else ""
        lines.append(
            f"| {row['engine']} | {row['fixture']} | {row['ocr_status']} "
            f"| {row['score']:.4f} | {row['elapsed_ms']:.2f} | {warn} |"
        )

    lines += ["", "## OCR Text Samples", ""]
    for row in rows:
        excerpt = row["ocr_text"].replace("\n", " / ")[:250]
        lines.append(f"- **{row['engine']}** / `{row['fixture']}`: `{excerpt}`")

    lines += [
        "", "## Raw JSON", "", "```json",
        json.dumps({"ground_truth": ground_truth, "environment": env, "rows": rows},
                   ensure_ascii=False, indent=2),
        "```",
    ]
    return "\n".join(lines) + "\n"


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> int:
    print("Extracting ground truth from native PDF …")
    ground_truth = _extract_ground_truth(ANALYST_PDF)
    print(f"  {len(ground_truth)} chars extracted")

    print("Building image-based fixtures …")
    png_bytes = ANALYST_PNG.read_bytes()
    fixtures = _build_fixtures(png_bytes)

    specs = [
        EngineSpec("rapidocr-local-cpu", "rapidocr", "rapidocr:ch_sim+en", ("ch_sim", "en"), "cpu"),
        EngineSpec("tesseract-local-cpu", "tesseract", "tesseract:ch_sim+en", ("ch_sim", "en"), "cpu"),
    ]
    if sys.platform == "darwin":
        specs.append(
            EngineSpec("rapidocr-macos-coreml", "rapidocr", "rapidocr:ch_sim+en", ("ch_sim", "en"), "coreml")
        )

    workdir = Path(tempfile.mkdtemp(prefix="realworld-ocr-bench-"))
    rows: list[dict] = []

    for spec in specs:
        os.environ.update({
            "OCR_PROVIDER": spec.provider,
            "OCR_MODEL": spec.model,
            "OCR_LANGUAGES": ",".join(spec.languages),
            "OCR_DEVICE": spec.device,
            "OCR_MODEL_STORAGE_DIR": str(workdir / f"models-{spec.provider}"),
        })
        try:
            runtime = OcrRuntime(Settings.from_env())
            provider = LocalRuntimeOcrProvider(spec, runtime)
            _run(provider, fixtures[0], spec.languages, ground_truth)  # warmup
            for fixture in fixtures:
                row = _run(provider, fixture, spec.languages, ground_truth)
                row["engine"] = spec.name
                rows.append(row)
                icon = "✓" if row["score"] > 0.5 else ("~" if row["score"] > 0 else "✗")
                print(f"  {icon} {spec.name}/{fixture.name}: score={row['score']:.3f}  {row['elapsed_ms']:.0f}ms")
        except Exception as exc:
            print(f"  ENGINE FAILED {spec.name}: {exc}")
            for fixture in fixtures:
                rows.append({
                    "engine": spec.name,
                    "fixture": fixture.name,
                    "ocr_status": "failed",
                    "elapsed_ms": 0.0,
                    "score": 0.0,
                    "ocr_text": "",
                    "warnings": [f"{type(exc).__name__}: {exc}"],
                })

    import onnxruntime as ort
    env = {
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "host": os.uname().nodename,
        "python": sys.version.split()[0],
        "onnxruntime_providers": ort.get_available_providers(),
    }

    output_path = ROOT / "docs" / "OCR_PROVIDER_REALWORLD_BENCHMARK.md"
    output_path.write_text(_markdown_report(ground_truth, rows, env), encoding="utf-8")
    print(f"\nWrote report to {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
